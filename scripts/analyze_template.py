from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation('Hendrix Modern Dark · SlidesMania.pptx')
print(f'Slide width: {prs.slide_width}, height: {prs.slide_height}')
print(f'Slide width (inches): {prs.slide_width/914400}, height (inches): {prs.slide_height/914400}')
print(f'Total slides: {len(prs.slides)}')
print(f'Total layouts: {len(prs.slide_layouts)}')
print()

for i, slide in enumerate(prs.slides):
    if i >= 15:
        break
    print(f'=== Slide {i+1} (layout: {slide.slide_layout.name}) ===')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = p.text.strip()
                if t:
                    font_info = ""
                    for run in p.runs:
                        try:
                            if run.font.size:
                                font_info += f" size={run.font.size/12700:.1f}pt"
                        except:
                            pass
                        try:
                            font_info += f" color={run.font.color.rgb}"
                        except:
                            pass
                        try:
                            if run.font.name:
                                font_info += f" font={run.font.name}"
                        except:
                            pass
                        try:
                            if run.font.bold:
                                font_info += " BOLD"
                        except:
                            pass
                        break
                    print(f'  TEXT: "{t[:120]}"{font_info}')
    print()

# Also check the slide master background
print("=== Layout names ===")
for i, layout in enumerate(prs.slide_layouts):
    print(f'  Layout {i}: {layout.name}')
